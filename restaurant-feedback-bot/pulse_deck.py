# -*- coding: utf-8 -*-
"""
PulseTeam — экспорт в .pptx (упрощённая вёрстка: шаблоны PowerPoint).

Для презентации «как лендинг» откройте pitch.html в браузере — тот же визуальный язык
(Mulish, цвета, карточки, сетка). PPTX из этого скрипта — запасной вариант для тех,
кто требует именно файл PowerPoint / WPS.
"""
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0xFA, 0xF7, 0xF3)
INK = RGBColor(0x1C, 0x19, 0x17)
MUTED = RGBColor(0x57, 0x53, 0x4E)
ACCENT = RGBColor(0xFF, 0x5A, 0x1F)
CALM = RGBColor(0x0D, 0x94, 0x88)


def set_slide_bg(slide, rgb):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb


def body_placeholder(slide):
    for shape in slide.placeholders:
        if shape.is_placeholder and shape.placeholder_format.idx == 1:
            return shape
    return None


def accent_top_strip(slide, prs):
    """Тонкая полоса у верхнего края — не пересекает заголовки слайдов."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Emu(40000)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = ACCENT
    shp.line.fill.background()


def accent_divider(slide, prs):
    """Линия-подчёркивание под зоной заголовка (ниже типового title)."""
    margin = Inches(0.56)
    top = Inches(1.72)
    width = prs.slide_width - 2 * margin
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, margin, top, width, Emu(32000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = ACCENT
    shp.line.fill.background()


def style_title(slide, size_pt=34):
    tf = slide.shapes.title.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_bottom = Pt(6)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(size_pt)
        p.font.bold = True
        p.font.color.rgb = INK


def style_title_slide(s0):
    tf = s0.shapes.title.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT
    sub = s0.placeholders[1]
    stf = sub.text_frame
    stf.word_wrap = True
    for p in stf.paragraphs:
        p.alignment = PP_ALIGN.LEFT


def fill_body(slide, lines, lead=None):
    ph = body_placeholder(slide)
    if not ph:
        return
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Emu(114300)
    tf.margin_right = Emu(114300)
    tf.margin_top = Pt(4)
    first = True
    if lead:
        p = tf.paragraphs[0]
        p.text = lead
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = CALM
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        first = False
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        if not line.strip():
            p.space_after = Pt(5)
            continue
        s = line.strip()
        is_dash = s.startswith("—") or s.startswith("–")
        is_subhead = s.endswith(":") and len(s) < 80 and not is_dash
        if is_subhead:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = CALM
            p.space_before = Pt(10)
            p.space_after = Pt(4)
        elif is_dash:
            p.font.size = Pt(15)
            p.font.color.rgb = MUTED
            p.space_after = Pt(5)
        else:
            p.font.size = Pt(16)
            p.font.color.rgb = INK
            p.space_after = Pt(6)


def add_slide(prs, title, lines, lead=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_bg(slide, BG)
    slide.shapes.title.text = title
    style_title(slide)
    accent_divider(slide, prs)
    fill_body(slide, lines, lead)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 — Титул
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(s0, BG)
    accent_top_strip(s0, prs)
    s0.shapes.title.text = "PulseTeam"
    tf0 = s0.shapes.title.text_frame
    tf0.paragraphs[0].font.size = Pt(54)
    tf0.paragraphs[0].font.bold = True
    tf0.paragraphs[0].font.color.rgb = ACCENT
    tf0.paragraphs[0].alignment = PP_ALIGN.LEFT
    sub = s0.placeholders[1]
    stf = sub.text_frame
    stf.clear()
    p0 = stf.paragraphs[0]
    p0.text = (
        "Система аналитики рабочей среды и настроения персонала для ресторанов и сетей HoReCa."
    )
    p0.font.size = Pt(22)
    p0.font.color.rgb = INK
    p0.alignment = PP_ALIGN.LEFT
    p1 = stf.add_paragraph()
    p1.text = (
        "Сигналы о выгорании, конфликтах и напряжённых сменах — до удара по сервису и выручке. "
        "Telegram-first, без отдельного приложения для линии."
    )
    p1.font.size = Pt(17)
    p1.font.color.rgb = MUTED
    p1.space_before = Pt(12)
    p1.alignment = PP_ALIGN.LEFT
    p2 = stf.add_paragraph()
    p2.text = "AI-инсайты · пилот и первые данные за короткий цикл"
    p2.font.size = Pt(16)
    p2.font.color.rgb = CALM
    p2.space_before = Pt(14)
    p2.alignment = PP_ALIGN.LEFT
    style_title_slide(s0)

    slides = [
        (
            "Проблема",
            None,
            [
                "Команда «ломается» раньше, чем это видно в цифрах смены: текучка, выгорание, токсичные паттерны общения, потеря мотивации.",
                "Сервис и выручка уже просели, когда руководство начинает разбираться.",
                "",
                "Типичные сигналы, которые замечают слишком поздно:",
                "— сотрудник выгорает и уходит",
                "— менеджер системно давит на смену",
                "— конфликты дозревают до кассы и отзывов",
                "— качество сервиса плавает без понятной причины",
                "",
                "Классический HR и редкие опросы не дают ежедневной картины по сменам — только постфактум.",
            ],
        ),
        (
            "Рынок и окно",
            None,
            [
                "В HoReCa нет лёгкого инструмента ежедневной аналитики атмосферы и состояния команды «в операционке».",
                "",
                "Тяжёлые HRM/корпоративные платформы плохо приживаются у линейного персонала; разовые опросы не ловят динамику смен.",
                "",
                "Запрос рынка:",
                "— быстро, с телефона, каждый день",
                "— минимум трения для сотрудника",
                "— понятные инсайты для владельца и управляющего",
                "— масштабирование на сеть без «второго проекта внедрения»",
            ],
        ),
        (
            "Решение PulseTeam",
            None,
            [
                "Короткий daily check-in в Telegram + аналитический слой для руководства.",
                "",
                "Что видим из данных:",
                "— настроение и напряжение команды по сменам",
                "— качество смен и повторяющиеся боли",
                "— вовлечённость и динамика по менеджерам и точкам",
                "",
                "Руководитель получает дашборд, алерты и AI-сводки — картина по сети, а не разрозненные отзывы.",
            ],
        ),
        (
            "Продукт: ценность",
            None,
            [
                "Для сотрудников:",
                "— быстрые pulse-check и анонимная обратная связь",
                "— Telegram без установки отдельного приложения",
                "— меньше 10 секунд на ответ, без «обучения системе»",
                "",
                "Для управляющих:",
                "— дашборд по сменам, проблемные сигналы, weekly-отчёты",
                "— heatmap настроения и сравнение филиалов",
                "",
                "Для сети:",
                "— multi-location аналитика, роли, единая модель данных по ресторанам",
            ],
        ),
        (
            "Как это работает",
            None,
            [
                "1. Подключаем ресторан: структура команды и роли.",
                "2. Сотрудники проходят check-in в Telegram (5–10 секунд до/после смены).",
                "3. Платформа выявляет негативные тренды и «узкие места».",
                "4. Руководитель получает аналитику: дашборд, отчёты, сигналы о рисках.",
            ],
        ),
        (
            "Инсайты для руководства",
            None,
            [
                "На одном экране — не сырые ответы, а смысл:",
                "— проблемные смены и просадки настроения",
                "— риски выгорания и снижение вовлечённости",
                "— зоны ответственности менеджеров по обратной связи",
                "— heatmap по точкам сети",
                "— AI summary по неделе вместо ручного сбора «что у нас в команде»",
            ],
        ),
        (
            "Аудитория и вовлечённость",
            None,
            [
                "Сегменты:",
                "— рестораны и сети, кофейни, бары, dark kitchen, delivery-команды",
                "",
                "Портрет покупателя:",
                "— владелец: текучка, стабильность команды, сервис",
                "— управляющий: атмосфера смен, прозрачность, ранние конфликты",
                "— HR / операционный директор: вовлечённость и динамика персонала",
                "",
                "Почему отвечают сотрудники:",
                "— привычный Telegram, мало вопросов, можно анонимно",
                "— не конкурирует с тяжёлыми HR-системами",
            ],
        ),
        (
            "Экономика: дешевле найма",
            None,
            [
                "Отраслевой ориентир по HoReCa (линия): годовая текучка нередко 60–100%+ — большая доля смены «обновляется» каждый год.",
                "",
                "Модельная стоимость одной замены в РФ (подбор, онбординг, просадка скорости, срывы смен): консервативно ≈40–120 тыс. ₽ на человека.",
                "",
                "Пример: точка ~25 человек, текучка 70% → ~17–18 уходов в год; при ~50 тыс. ₽ на замену — порядка 0,8–0,9 млн ₽/год только на цикл замены (без упущенной выручки).",
                "",
                "Подписка 4,9–9,9 тыс. ₽/мес ≈ 60–120 тыс. ₽/год — дешевле одного цикла найма. Если ранние сигналы помогают удержать 2–3 человека в год, экономия на заменах часто в несколько раз выше стоимости сервиса.",
                "Цифры — иллюстрация; покупатель подставляет свою численность и текучку.",
            ],
        ),
        (
            "Монетизация и масштаб",
            None,
            [
                "Модель: SaaS на ресторан + разовое подключение сети.",
                "",
                "Тарифы (ориентиры с лендинга):",
                "— подключение сети: от 30 000 ₽",
                "— до 10 сотрудников: 2 990 ₽/мес",
                "— 11–25 сотрудников: 6 990 ₽/мес",
                "— 26–40 сотрудников: 8 990 ₽/мес",
                "— 40+: индивидуально",
                "",
                "Иллюстрация верхней оценки MRR: 200 точек × ~8 000 ₽ ≈ 1,6 млн ₽/мес; 1 000 точек — порядка 8 млн ₽/мес при сопоставимом ARPA.",
                "Высокая маржинальность за счёт SaaS и относительно низкой инфраструктурной нагрузки.",
            ],
        ),
        (
            "Тайминг и отстройка",
            None,
            [
                "Почему сейчас:",
                "— кадровый дефицит и рост текучки в HoReCa",
                "— Telegram уже стандарт коммуникации на сменах",
                "— AI удешевляет разбор сигналов и отчёты",
                "— фокус смещается к удержанию людей, не только к найму",
                "",
                "Почему мы, а не «ещё один HR-SaaS»:",
                "— узкий фокус: смены, атмосфера, операционная аналитика",
                "— нативный Telegram и ежедневный сценарий использования",
                "— AI поверх регулярных коротких данных, а не разовые опросы",
            ],
        ),
        (
            "Видение",
            None,
            [
                "Эволюция: от регулярной обратной связи к predictive-слою — прогноз ухода, рисковые смены, корреляции настроения с операционными метриками.",
                "",
                "Следующий горизонт — AI-assistant для управленческих решений в ресторане (подсказки, а не «чёрный ящик»).",
            ],
        ),
        (
            "Контакты",
            None,
            [
                "Пилот, демо, тарифы и подключение точки или сети.",
                "",
                "Заявки и диалог в Telegram: https://t.me/plsbuyit",
                "",
                "Бот для ежедневной обратной связи с линии (оценка смены в личке): https://t.me/smena_feedback_bot",
            ],
        ),
    ]

    for title, lead, lines in slides:
        add_slide(prs, title, lines, lead)

    script_dir = Path(__file__).resolve().parent
    out_path = script_dir / "PulseTeam_Pitch_Deck.pptx"
    print("Папка скрипта:", script_dir)
    print("Запись файла:", out_path)

    try:
        prs.save(str(out_path))
    except (PermissionError, OSError) as e:
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        alt = script_dir / f"PulseTeam_Pitch_Deck_{stamp}.pptx"
        print(
            "Не удалось перезаписать основной файл (закройте его в WPS / PowerPoint).",
            "Ошибка:",
            e,
            sep="\n",
        )
        print("Сохраняю копию:", alt)
        prs.save(str(alt))
        out_path = alt

    print("Готово. Откройте:", out_path)
    print("Слайдов (с титулом):", len(prs.slides))


if __name__ == "__main__":
    main()
